#!/usr/bin/env python3
"""
Build AI-Deployment-Framework-v4.2.pptx by loading v4.1 and appending 4 appendix
slides with generated capability graph images.
"""

import sys
from pathlib import Path

sys.path.insert(0, "/Users/bletort/.cursor/skills/pptx-builder/scripts")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx_helpers import (
    THEMES,
    PP_ALIGN,
    add_header_slide,
    add_footer,
    add_callout_box,
    _textbox,
    _set_text,
)

BASE = Path("/Users/bletort/Library/CloudStorage/OneDrive-DigitalRealty/Shared")
V41 = BASE / "AI-Deployment-Framework-v4.1.pptx"
V42 = BASE / "AI-Deployment-Framework-v4.2.pptx"
ASSETS = Path("/Users/bletort/Data/MemoryOS/assets")

FOOTER = "Digital Realty  |  AI Deployment Framework"

IMG_LEFT = Inches(0.4)
IMG_TOP = Inches(1.0)
IMG_WIDTH = Inches(12.5)
IMG_HEIGHT = Inches(5.3)


def add_image_slide(
    prs,
    colors: dict,
    title: str,
    image_path: Path,
    callout_text: str | None = None,
    callout_style: str = "info",
) -> None:
    slide = add_header_slide(prs, title, colors)
    slide.shapes.add_picture(str(image_path), IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT)
    if callout_text:
        add_callout_box(
            slide,
            callout_text,
            Inches(0.5),
            Inches(6.4),
            Inches(12.3),
            Inches(0.6),
            colors,
            style=callout_style,
        )
    add_footer(slide, FOOTER, colors)


def main():
    if not V41.exists():
        raise SystemExit(f"Source not found: {V41}")

    prs = Presentation(str(V41))
    colors = THEMES["light"]

    add_image_slide(
        prs,
        colors,
        "Appendix: Global Business Process Capability Graph",
        ASSETS / "slide16_capability_graph_all_in_one.png",
        callout_text="Intents + Agents augment workers: the capability graph executes work end-to-end across BUs.",
    )

    add_image_slide(
        prs,
        colors,
        "Appendix: Example Path — Contract Renewal Intent",
        ASSETS / "slide17_legal_renewal_path.png",
        callout_text='Intent: "Renew this contract" — Agents chain capabilities across Salesforce, Carma, and Legal.',
    )

    add_image_slide(
        prs,
        colors,
        "Appendix: Example Path — GTM Outreach Intent",
        ASSETS / "slide18_gtm_outreach_path.png",
        callout_text='Intent: "Generate outreach sequence" — Same pattern: intent drives the graph, agents execute, humans review.',
    )

    add_image_slide(
        prs,
        colors,
        "Appendix: Manual vs Agentic Execution",
        ASSETS / "slide19_manual_vs_agentic.png",
        callout_text="Same business processes — different execution model. Agents augment workers.",
        callout_style="success",
    )

    prs.save(str(V42))
    print(f"Saved: {V42}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
