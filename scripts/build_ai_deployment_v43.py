#!/usr/bin/env python3
"""
Build AI-Deployment-Framework-v4.3.pptx from v4.1 base with embedded
capability graph images as appendix slides.
"""

import sys
from pathlib import Path

sys.path.insert(0, "/Users/bletort/.cursor/skills/pptx-builder/scripts")
from pptx import Presentation
from pptx.util import Inches
from pptx_helpers import (
    THEMES,
    add_header_slide,
    add_footer,
    add_callout_box,
)

BASE = Path("/Users/bletort/Library/CloudStorage/OneDrive-DigitalRealty/Shared")
V41 = BASE / "AI-Deployment-Framework-v4.1.pptx"
V43 = BASE / "AI-Deployment-Framework-v4.3.pptx"
ASSETS = Path("/Users/bletort/Data/MemoryOS/assets/v43")

FOOTER = "Digital Realty  |  AI Deployment Framework"

IMG_LEFT = Inches(0.4)
IMG_TOP = Inches(1.0)
IMG_WIDTH = Inches(12.5)
IMG_HEIGHT = Inches(5.3)

SLIDES = [
    {
        "title": "Appendix: Global Business Process Capability Graph",
        "image": "slide16_capability_graph_v2.png",
        "callout": "Intents + Agents augment workers: the capability graph executes work end-to-end across BUs.",
        "style": "info",
    },
    {
        "title": "Appendix: Example Path — Contract Renewal Intent",
        "image": "slide17_legal_renewal_v2.png",
        "callout": 'Intent: "Renew this contract" — Agents chain capabilities across Salesforce, Carma, and Legal.',
        "style": "info",
    },
    {
        "title": "Appendix: Example Path — GTM Outreach Intent",
        "image": "slide18_gtm_outreach_v2.png",
        "callout": 'Intent: "Generate outreach sequence" — Same pattern: intent drives the graph, agents execute, humans review.',
        "style": "info",
    },
    {
        "title": "Appendix: Manual vs Agentic Execution",
        "image": "slide19_manual_vs_agentic.png",
        "callout": "Same business processes — different execution model. Agents augment workers.",
        "style": "success",
    },
]


def main():
    if not V41.exists():
        raise SystemExit(f"Source not found: {V41}")

    for s in SLIDES:
        img = ASSETS / s["image"]
        if not img.exists():
            raise SystemExit(f"Image not found: {img}")

    prs = Presentation(str(V41))
    colors = THEMES["light"]

    for s in SLIDES:
        slide = add_header_slide(prs, s["title"], colors)
        slide.shapes.add_picture(
            str(ASSETS / s["image"]), IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT
        )
        add_callout_box(
            slide,
            s["callout"],
            Inches(0.5),
            Inches(6.4),
            Inches(12.3),
            Inches(0.6),
            colors,
            style=s["style"],
        )
        add_footer(slide, FOOTER, colors)

    prs.save(str(V43))
    print(f"Saved: {V43}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"File size: {V43.stat().st_size / 1024 / 1024:.1f} MB")


if __name__ == "__main__":
    main()
