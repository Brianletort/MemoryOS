"""File upload handling -- extract text from any file type."""

from __future__ import annotations

import csv
import io
import logging
import uuid
from pathlib import Path
from typing import Any

logger = logging.getLogger("memoryos.chat.files")

REPO_DIR = Path(__file__).resolve().parent.parent.parent
UPLOAD_DIR = REPO_DIR / "data" / "uploads"
GENERATED_DIR = REPO_DIR / "data" / "generated"

MAX_FILE_SIZE = 100 * 1024 * 1024  # 100 MB

TEXT_EXTENSIONS = {
    ".txt", ".md", ".json", ".yaml", ".yml", ".xml", ".html", ".htm",
    ".css", ".js", ".ts", ".jsx", ".tsx", ".py", ".rb", ".go", ".rs",
    ".java", ".c", ".cpp", ".h", ".hpp", ".cs", ".swift", ".kt",
    ".sql", ".sh", ".bash", ".zsh", ".ps1", ".bat", ".cmd",
    ".log", ".ini", ".cfg", ".conf", ".toml", ".env",
    ".r", ".m", ".pl", ".lua", ".dart", ".scala",
    ".rtf", ".tex", ".bib",
}


def _ensure_dirs() -> None:
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)


def save_upload(filename: str, content: bytes) -> dict[str, Any]:
    """Save an uploaded file and return metadata with extracted text."""
    _ensure_dirs()

    if len(content) > MAX_FILE_SIZE:
        raise ValueError(f"File too large ({len(content)} bytes). Max: {MAX_FILE_SIZE} bytes")

    ext = Path(filename).suffix.lower()
    file_id = uuid.uuid4().hex[:12]
    safe_name = f"{file_id}_{Path(filename).stem[:50]}{ext}"
    dest = UPLOAD_DIR / safe_name
    dest.write_bytes(content)

    extracted = extract_text(dest)

    return {
        "file_id": file_id,
        "filename": filename,
        "stored_as": safe_name,
        "size": len(content),
        "extension": ext,
        "text": extracted,
        "text_length": len(extracted),
    }


def extract_text(file_path: Path) -> str:
    """Extract readable text from a file based on its extension."""
    ext = file_path.suffix.lower()

    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        elif ext in (".docx", ".doc"):
            return _extract_docx(file_path)
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(file_path)
        elif ext == ".csv":
            return _extract_csv(file_path)
        elif ext in (".xlsx", ".xls"):
            return _extract_excel(file_path)
        elif ext in TEXT_EXTENSIONS:
            return _extract_text(file_path)
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg", ".bmp", ".ico"):
            return f"[Image file: {file_path.name}]"
        elif ext in (".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac"):
            return f"[Audio file: {file_path.name}]"
        elif ext in (".mp4", ".mov", ".avi", ".mkv", ".webm"):
            return f"[Video file: {file_path.name}]"
        elif ext == ".zip":
            return _extract_zip_listing(file_path)
        else:
            return _try_text_fallback(file_path)
    except Exception as exc:
        logger.warning("Text extraction failed for %s: %s", file_path, exc)
        return f"[Error extracting text from {file_path.name}: {exc}]"


def _extract_text(path: Path) -> str:
    content = path.read_text(encoding="utf-8", errors="replace")
    if len(content) > 100_000:
        content = content[:100_000] + f"\n\n... [truncated, file is {len(content)} chars]"
    return content


def _try_text_fallback(path: Path) -> str:
    """Try reading as text; if it fails, report as binary."""
    try:
        content = path.read_text(encoding="utf-8", errors="strict")
        if len(content) > 100_000:
            content = content[:100_000] + "\n... [truncated]"
        return content
    except (UnicodeDecodeError, ValueError):
        return f"[Binary file: {path.name}, {path.stat().st_size} bytes]"


def _extract_pdf(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                parts.append(f"--- Page {i + 1} ---\n{text}")
    return "\n\n".join(parts) if parts else "[PDF contained no extractable text]"


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(parts) if parts else "[DOCX contained no text]"


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
    return "\n\n".join(parts) if parts else "[PPTX contained no text]"


def _extract_csv(path: Path) -> str:
    text = path.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return "[Empty CSV]"

    header = rows[0]
    preview_rows = rows[1:51]
    lines = [", ".join(header)]
    for row in preview_rows:
        lines.append(", ".join(row))
    if len(rows) > 51:
        lines.append(f"... ({len(rows) - 1} total rows)")
    return "\n".join(lines)


def _extract_excel(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(max_row=51, values_only=True))
        if not rows:
            continue
        lines = [f"## Sheet: {sheet_name}"]
        for row in rows:
            lines.append(", ".join(str(c) if c is not None else "" for c in row))
        parts.append("\n".join(lines))
    wb.close()
    return "\n\n".join(parts) if parts else "[Excel file contained no data]"


def _extract_zip_listing(path: Path) -> str:
    import zipfile

    try:
        with zipfile.ZipFile(path, "r") as zf:
            names = zf.namelist()
            lines = [f"ZIP archive: {len(names)} files"]
            for name in names[:100]:
                info = zf.getinfo(name)
                lines.append(f"  {name} ({info.file_size} bytes)")
            if len(names) > 100:
                lines.append(f"  ... and {len(names) - 100} more")
            return "\n".join(lines)
    except Exception:
        return f"[ZIP file: {path.name}]"


def get_upload_path(stored_name: str) -> Path | None:
    """Return the full path for a stored upload, or None if not found."""
    path = UPLOAD_DIR / stored_name
    return path if path.is_file() else None


def get_generated_path(filename: str) -> Path | None:
    """Return the full path for a generated file, or None if not found."""
    path = GENERATED_DIR / filename
    return path if path.is_file() else None
