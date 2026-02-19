#!/usr/bin/env python3
"""OneDrive local sync folder -> Obsidian Markdown extractor.

Scans the locally-synced OneDrive folder for Office documents and converts
them to searchable Markdown files in the Obsidian vault:
  - .docx -> pandoc (requires `brew install pandoc`)
  - .pptx -> python-pptx (slide titles + bullet text)
  - .pdf  -> pdfplumber (text extraction, best-effort)
  - .xlsx -> openpyxl (sheet names + first rows as tables)

Tracks file modification times in state.json for incremental updates.
"""

from __future__ import annotations

import argparse
import logging
import os
import subprocess
import sys
from datetime import datetime
from pathlib import Path
from typing import Any

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))

from src.common.config import load_config, resolve_output_dir, setup_logging
from src.common.markdown import sanitize_filename, write_markdown, yaml_frontmatter
from src.common.state import get_cursor, load_state, save_state, set_cursor

logger = logging.getLogger("memoryos.onedrive")


# ── Converters ───────────────────────────────────────────────────────────────

def convert_docx(path: Path) -> str | None:
    """Convert .docx to Markdown using pandoc."""
    try:
        result = subprocess.run(
            ["pandoc", "-f", "docx", "-t", "markdown", "--wrap=none", str(path)],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
        logger.warning("pandoc failed for %s: %s", path, result.stderr[:200])
        return None
    except FileNotFoundError:
        logger.error("pandoc not found. Install with: brew install pandoc")
        return None
    except subprocess.TimeoutExpired:
        logger.warning("pandoc timed out for %s", path)
        return None


def convert_pptx(path: Path) -> str | None:
    """Convert .pptx to Markdown using python-pptx."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        logger.error("python-pptx not installed. Run: pip install python-pptx")
        return None

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        logger.warning("Failed to open %s: %s", path, exc)
        return None

    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body_lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False:
                        title = text
                    else:
                        level = para.level or 0
                        indent = "  " * level
                        body_lines.append(f"{indent}- {text}")

        parts.append(f"## Slide {i}" + (f": {title}" if title else ""))
        parts.append("")
        if body_lines:
            parts.extend(body_lines)
            parts.append("")

    return "\n".join(parts) if parts else None


def convert_pdf(path: Path) -> str | None:
    """Extract text from PDF using pdfplumber."""
    try:
        import pdfplumber  # type: ignore[import-untyped]
    except ImportError:
        logger.error("pdfplumber not installed. Run: pip install pdfplumber")
        return None

    try:
        parts: list[str] = []
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    parts.append(f"## Page {i}")
                    parts.append("")
                    parts.append(text.strip())
                    parts.append("")
        return "\n".join(parts) if parts else None
    except Exception as exc:
        logger.warning("Failed to extract PDF %s: %s", path, exc)
        return None


def convert_xlsx(path: Path) -> str | None:
    """Extract content from Excel using openpyxl."""
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError:
        logger.error("openpyxl not installed. Run: pip install openpyxl")
        return None

    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        logger.warning("Failed to open %s: %s", path, exc)
        return None

    parts: list[str] = []
    max_rows = 100  # Limit rows per sheet

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## Sheet: {sheet_name}")
        parts.append("")

        rows_written = 0
        header: list[str] = []
        for row in ws.iter_rows(max_col=20, values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if not any(cells):
                continue
            if rows_written == 0:
                header = cells
                parts.append("| " + " | ".join(cells) + " |")
                parts.append("| " + " | ".join(["---"] * len(cells)) + " |")
            else:
                parts.append("| " + " | ".join(cells) + " |")
            rows_written += 1
            if rows_written >= max_rows:
                parts.append(f"*... ({ws.max_row - max_rows} more rows)*")
                break

        parts.append("")

    wb.close()
    return "\n".join(parts) if parts else None


CONVERTERS = {
    ".docx": convert_docx,
    ".pptx": convert_pptx,
    ".pdf": convert_pdf,
    ".xlsx": convert_xlsx,
}


# ── Main extraction logic ───────────────────────────────────────────────────

def run(cfg: dict[str, Any], *, dry_run: bool = False) -> None:
    """Run the OneDrive extractor."""
    sync_dir = Path(cfg["onedrive"]["sync_dir"])
    if not sync_dir.is_dir():
        logger.error("OneDrive sync dir not found: %s", sync_dir)
        return

    od_settings = cfg.get("onedrive_settings", {})
    extensions = set(cfg["onedrive"].get("extensions", [".docx", ".pptx", ".pdf", ".xlsx"]))
    max_size_mb = od_settings.get("max_file_size_mb", 50)
    max_size_bytes = max_size_mb * 1024 * 1024
    max_retries: int = od_settings.get("max_conversion_retries", 3)

    state_path = cfg["state_file"]
    state = load_state(state_path)

    file_mtimes: dict[str, float] = get_cursor(state, "onedrive", "file_mtimes", {})
    failed_files: dict[str, dict[str, Any]] = get_cursor(state, "onedrive", "failed_files", {})

    slides_dir = resolve_output_dir(cfg, "slides")
    knowledge_dir = resolve_output_dir(cfg, "knowledge")

    files_to_process: list[Path] = []
    skipped_failures = 0
    for ext in extensions:
        for fpath in sync_dir.rglob(f"*{ext}"):
            if fpath.name.startswith("~$"):
                continue
            if fpath.stat().st_size > max_size_bytes:
                logger.debug("Skipping large file: %s (%d MB)", fpath, fpath.stat().st_size // (1024*1024))
                continue
            current_mtime = fpath.stat().st_mtime
            cached_mtime = file_mtimes.get(str(fpath), 0)

            fail_info = failed_files.get(str(fpath))
            if (
                fail_info
                and fail_info.get("mtime") == current_mtime
                and fail_info.get("count", 0) >= max_retries
            ):
                skipped_failures += 1
                continue

            if current_mtime > cached_mtime:
                files_to_process.append(fpath)

    if skipped_failures:
        logger.info(
            "Skipping %d files with repeated failures (will retry if files change)",
            skipped_failures,
        )

    if not files_to_process:
        logger.info("No new or modified OneDrive files")
        if not dry_run and skipped_failures:
            set_cursor(state, "onedrive", "failed_files", failed_files)
            save_state(state_path, state)
        return

    logger.info("Processing %d OneDrive files", len(files_to_process))

    converted = 0
    failed = 0

    for fpath in files_to_process:
        ext = fpath.suffix.lower()
        converter = CONVERTERS.get(ext)
        if not converter:
            logger.debug("No converter for %s", ext)
            continue

        target_dir = slides_dir if ext == ".pptx" else knowledge_dir
        rel = fpath.relative_to(sync_dir)
        out_name = sanitize_filename(fpath.stem) + ".md"
        out_path = target_dir / rel.parent / out_name

        logger.info("Converting: %s -> %s", fpath.name, out_path)

        if dry_run:
            logger.info("DRY RUN: Would convert %s", fpath)
            continue

        body = converter(fpath)
        if body is None:
            prev = failed_files.get(str(fpath), {})
            count = prev.get("count", 0) + 1
            failed_files[str(fpath)] = {
                "count": count,
                "mtime": fpath.stat().st_mtime,
                "last_error": "conversion produced no content",
            }
            if count >= max_retries:
                logger.warning(
                    "Skipping %s (failed %d times, will retry if file changes)",
                    fpath.name, count,
                )
            else:
                logger.warning("Conversion failed for %s (attempt %d/%d)", fpath.name, count, max_retries)
            failed += 1
            continue

        failed_files.pop(str(fpath), None)

        meta = {
            "source": "onedrive",
            "original_path": str(fpath),
            "converted_at": datetime.now().isoformat(timespec="seconds"),
            "type": ext.lstrip("."),
            "file_size": fpath.stat().st_size,
        }
        content = yaml_frontmatter(meta) + "\n\n"
        content += f"# {fpath.stem}\n\n"
        content += body

        write_markdown(out_path, content)
        converted += 1

        file_mtimes[str(fpath)] = fpath.stat().st_mtime

    if not dry_run:
        set_cursor(state, "onedrive", "file_mtimes", file_mtimes)
        set_cursor(state, "onedrive", "failed_files", failed_files)
        save_state(state_path, state)

    logger.info("OneDrive extraction complete: %d converted, %d failed", converted, failed)


# ── CLI ──────────────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(description="OneDrive -> Obsidian Markdown")
    parser.add_argument("--config", help="Path to config.yaml")
    parser.add_argument("--dry-run", action="store_true", help="Preview without writing")
    parser.add_argument("--reset", action="store_true", help="Reset to re-process all files")
    args = parser.parse_args()

    cfg = load_config(args.config)
    setup_logging(cfg)

    if args.reset:
        state = load_state(cfg["state_file"])
        set_cursor(state, "onedrive", "file_mtimes", {})
        set_cursor(state, "onedrive", "failed_files", {})
        save_state(cfg["state_file"], state)
        logger.info("Reset OneDrive file tracking state")

    run(cfg, dry_run=args.dry_run)


if __name__ == "__main__":
    main()
